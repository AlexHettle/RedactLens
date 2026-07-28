import io
import random
import zipfile
from pathlib import Path

import pytest
from docx import Document
from odf.opendocument import load as load_odf
from odf.text import P
from openpyxl import load_workbook
from pptx import Presentation
from redactlens_core.extractors import extract_document
from redactlens_core.validators import luhn_valid

from generate_corpus import (
    CALIBRATION_SEED,
    CORPUS_VERSION,
    FABRICATION_POLICY,
    HOLDOUT_SEED,
    STRUCTURE_SIGNATURE_VERSION,
    TEMPLATE_FAMILIES,
    CorpusBundle,
    CorpusDocument,
    Plant,
    _structured_topology_signature,
    fake_invalid_card_number,
    fake_valid_card_number,
    generate_all,
    manifest_for,
    role_separation_evidence,
    validate_role_separation,
    write_all_corpora,
    write_bundle,
)


def _document_text(relative_path: str, content: bytes) -> str:
    suffix = Path(relative_path).suffix.lower()
    if suffix in {".docx", ".xlsx", ".pptx", ".odt"}:
        extracted = extract_document(suffix, content)
        assert extracted is not None
        return extracted.text
    return content.decode("utf-8")


def test_card_generators_agree_with_luhn_validator():
    rng = random.Random(99)
    for _ in range(20):
        assert luhn_valid(fake_valid_card_number(rng))
        assert not luhn_valid(fake_invalid_card_number(rng))


def test_corpora_are_deterministic_versioned_and_role_separated():
    first = generate_all()
    second = generate_all()

    assert first == second
    assert first["calibration"].seed == CALIBRATION_SEED
    assert first["holdout"].seed == HOLDOUT_SEED
    assert manifest_for(first["calibration"])["corpus_version"] == CORPUS_VERSION
    assert manifest_for(first["calibration"])["sha256"] != manifest_for(first["holdout"])["sha256"]
    assert CORPUS_VERSION == "3.0.0"

    calibration_manifest = manifest_for(first["calibration"])
    holdout_manifest = manifest_for(first["holdout"])
    assert calibration_manifest["template_family"] == TEMPLATE_FAMILIES["calibration"]
    assert holdout_manifest["template_family"] == TEMPLATE_FAMILIES["holdout"]
    assert calibration_manifest["template_family"] != holdout_manifest["template_family"]
    assert calibration_manifest["structure_signature_version"] == STRUCTURE_SIGNATURE_VERSION
    assert calibration_manifest["fabrication_policy"] == FABRICATION_POLICY


def test_every_document_path_and_case_id_is_unique_within_each_role():
    for bundle in generate_all().values():
        paths = [document.relative_path for document in bundle.documents]
        case_ids = [plant.case_id for plant in bundle.plants]
        assert len(paths) == len(set(paths))
        assert len(case_ids) == len(set(case_ids))


def test_every_plant_offset_matches_real_text_and_structured_extraction():
    for bundle in generate_all().values():
        documents = {
            document.relative_path: _document_text(document.relative_path, document.content)
            for document in bundle.documents
        }
        for plant in bundle.plants:
            text = documents[plant.file]
            assert 0 <= plant.start < plant.end <= len(text)
            value = text[plant.start : plant.end]
            assert value == value.strip()


def test_manifest_counts_all_generated_files_including_unlabeled_clean_files():
    for bundle in generate_all().values():
        manifest = manifest_for(bundle)
        labeled_files = {plant.file for plant in bundle.plants}
        assert manifest["document_count"] == len(bundle.documents)
        assert manifest["document_count"] > len(labeled_files)


def test_corpora_cover_required_hard_negatives_and_realistic_bundles():
    required_cases = {
        "negative-uuid",
        "negative-migration",
        "negative-lock-integrity",
        "negative-long-url",
        "negative-base64",
        "negative-test-password",
        "negative-test-aws",
        "negative-example-connection",
        "negative-placeholder-password",
        "negative-email-shaped-connection",
        "negative-public-email",
        "negative-ssn-shaped-id",
    }
    required_roots = {
        "python_service",
        "react_application",
        "infrastructure",
        "documentation",
        "office_documents",
    }
    for bundle in generate_all().values():
        case_ids = {plant.case_id for plant in bundle.plants}
        paths = {document.relative_path for document in bundle.documents}
        assert all(any(required in case_id for case_id in case_ids) for required in required_cases)
        assert required_roots <= {path.split("/", 1)[0] for path in paths}


def test_holdout_values_are_not_reused_from_calibration():
    bundles = generate_all()
    values_by_role = {}
    for role, bundle in bundles.items():
        texts = {
            document.relative_path: _document_text(document.relative_path, document.content)
            for document in bundle.documents
        }
        values_by_role[role] = {
            texts[plant.file][plant.start : plant.end]
            for plant in bundle.plants
            if plant.is_positive
        }
    assert values_by_role["calibration"].isdisjoint(values_by_role["holdout"])


def test_calibration_and_holdout_have_no_shared_structures_or_plant_contexts():
    bundles = generate_all()

    evidence = role_separation_evidence(bundles["calibration"], bundles["holdout"])

    assert evidence == {
        "signature_version": STRUCTURE_SIGNATURE_VERSION,
        "planted_value_overlaps": [],
        "plant_context_overlaps": [],
        "document_structure_overlaps": [],
        "structured_topology_overlaps": [],
    }


def test_role_separation_guard_rejects_a_relabelled_calibration_family():
    calibration = generate_all()["calibration"]
    relabelled = CorpusBundle(
        role="holdout",
        seed=HOLDOUT_SEED,
        documents=calibration.documents,
        plants=calibration.plants,
    )

    with pytest.raises(ValueError, match="independence guard failed"):
        validate_role_separation(calibration, relabelled)


def test_role_separation_signatures_ignore_detector_category_and_polarity_labels():
    calibration_text = b"prefix AAA suffix"
    holdout_text = b"prefix BBB suffix"
    calibration = CorpusBundle(
        role="calibration",
        seed=CALIBRATION_SEED,
        documents=(CorpusDocument("calibration/a.txt", calibration_text),),
        plants=(Plant("calibration/a.txt", 7, 10, "credential", True, "alpha", "a"),),
    )
    holdout = CorpusBundle(
        role="holdout",
        seed=HOLDOUT_SEED,
        documents=(CorpusDocument("holdout/b.txt", holdout_text),),
        plants=(Plant("holdout/b.txt", 7, 10, "personal_id", False, "beta", "b"),),
    )

    evidence = role_separation_evidence(calibration, holdout)

    assert evidence["plant_context_overlaps"]
    assert evidence["document_structure_overlaps"]
    with pytest.raises(ValueError, match="independence guard failed"):
        validate_role_separation(calibration, holdout)


def test_structured_topologies_are_role_distinct_and_guarded():
    bundles = generate_all()
    by_role = {}
    for role, bundle in bundles.items():
        by_role[role] = {
            Path(document.relative_path).suffix.lower(): _structured_topology_signature(document)
            for document in bundle.documents
            if Path(document.relative_path).suffix.lower() in {".docx", ".xlsx", ".pptx", ".odt"}
        }

    assert set(by_role["calibration"]) == {".docx", ".xlsx", ".pptx", ".odt"}
    assert set(by_role["holdout"]) == set(by_role["calibration"])
    for suffix, calibration_signature in by_role["calibration"].items():
        assert calibration_signature is not None
        assert calibration_signature != by_role["holdout"][suffix]

    calibration = bundles["calibration"]
    copied_document = next(
        document for document in calibration.documents if document.relative_path.endswith(".docx")
    )
    copied_topology = CorpusBundle(
        role="holdout",
        seed=HOLDOUT_SEED,
        documents=(copied_document,),
        plants=(),
    )
    evidence = role_separation_evidence(calibration, copied_topology)

    assert evidence["structured_topology_overlaps"] == [
        _structured_topology_signature(copied_document)
    ]
    with pytest.raises(ValueError, match="structured_topology_overlaps"):
        validate_role_separation(calibration, copied_topology)


def test_structured_topology_ignores_opaque_payload_bytes():
    def package(payload: bytes) -> bytes:
        output = io.BytesIO()
        with zipfile.ZipFile(output, "w") as archive:
            archive.writestr("x.xml", "<root><child /></root>")
            archive.writestr("media.bin", payload)
        return output.getvalue()

    first = CorpusDocument("first.docx", package(b"first payload"))
    second = CorpusDocument("second.docx", package(b"different payload"))

    assert _structured_topology_signature(first) == _structured_topology_signature(second)


def _planted_values(bundle):
    texts = {
        document.relative_path: _document_text(document.relative_path, document.content)
        for document in bundle.documents
    }
    return [(plant, texts[plant.file][plant.start : plant.end]) for plant in bundle.plants]


def test_identity_and_payment_fixtures_use_reserved_or_published_test_ranges():
    published_test_cards = {
        "4242424242424242",
        "5555555555554444",
        "378282246310005",
        "6011111111111117",
        "4000056655665556",
        "5200828282828210",
        "371449635398431",
        "6011000990139424",
        "30569309025904",
    }

    for role, bundle in generate_all().items():
        for plant, value in _planted_values(bundle):
            if plant.detector_id == "us_ssn":
                area, group, _ = (int(part) for part in value.split("-"))
                assert area in {0, 666} or 900 <= area <= 999 or group == 0
            elif plant.detector_id == "phone":
                _, exchange, subscriber = value.split("-")
                assert exchange == "555"
                assert 100 <= int(subscriber) <= 199
            elif plant.detector_id == "credit_card" and plant.is_positive:
                assert value in published_test_cards
                assert luhn_valid(value)
            elif plant.detector_id == "email":
                assert value.endswith("@fabricated.invalid")
            elif plant.detector_id == "connection_string":
                assert ".fabricated.invalid" in value

        # Role partitioning for the reserved SSN range makes accidental value
        # reuse impossible even before the structural guard runs.
        positive_ssns = [
            value
            for plant, value in _planted_values(bundle)
            if plant.is_positive and plant.detector_id == "us_ssn"
        ]
        expected_area = range(900, 950) if role == "calibration" else range(950, 1000)
        assert all(int(value[:3]) in expected_area for value in positive_ssns)


def test_structured_fixtures_are_valid_packages_with_fixed_zip_metadata():
    for bundle in generate_all().values():
        for document in bundle.documents:
            suffix = Path(document.relative_path).suffix.lower()
            if suffix not in {".docx", ".xlsx", ".pptx", ".odt"}:
                continue
            with zipfile.ZipFile(io.BytesIO(document.content)) as archive:
                names = archive.namelist()
                expected_order = sorted(names, key=lambda name: (name != "mimetype", name))
                assert names == expected_order
                for info in archive.infolist():
                    assert info.date_time == (1980, 1, 1, 0, 0, 0)
                    assert info.create_system == 0
                    assert info.external_attr == 0x20
                    assert info.compress_type == zipfile.ZIP_STORED
                    assert info.extra == b""
                    assert info.comment == b""

                if suffix == ".odt":
                    assert names[0] == "mimetype"
                    assert archive.read("mimetype") == b"application/vnd.oasis.opendocument.text"
                    assert "META-INF/manifest.xml" in names
                else:
                    assert "[Content_Types].xml" in names
                    assert "_rels/.rels" in names


def test_structured_fixtures_open_in_required_document_libraries():
    workbook_locations = {
        "calibration": ("Review", "B2"),
        "holdout": ("Casework", "D4"),
    }
    expected_fragments = {
        "calibration": {
            ".docx": "customer ssn",
            ".xlsx": "card on file",
            ".pptx": "personal contact email",
            ".odt": "call customer phone",
        },
        "holdout": {
            ".docx": "SSN",
            ".xlsx": "payment card",
            ".pptx": "email",
            ".odt": "mobile number",
        },
    }

    for role, bundle in generate_all().items():
        documents = {
            Path(document.relative_path).suffix.lower(): document
            for document in bundle.documents
            if Path(document.relative_path).suffix.lower() in {".docx", ".xlsx", ".pptx", ".odt"}
        }

        loaded_docx = Document(io.BytesIO(documents[".docx"].content))
        docx_text = [paragraph.text for paragraph in loaded_docx.paragraphs]
        docx_text.extend(
            cell.text for table in loaded_docx.tables for row in table.rows for cell in row.cells
        )
        assert expected_fragments[role][".docx"] in "\n".join(docx_text)

        loaded_xlsx = load_workbook(io.BytesIO(documents[".xlsx"].content), read_only=True)
        sheet_name, cell_reference = workbook_locations[role]
        try:
            assert (
                expected_fragments[role][".xlsx"] in loaded_xlsx[sheet_name][cell_reference].value
            )
        finally:
            loaded_xlsx.close()

        loaded_pptx = Presentation(io.BytesIO(documents[".pptx"].content))
        slide_text = "\n".join(
            shape.text for shape in loaded_pptx.slides[0].shapes if shape.has_text_frame
        )
        assert expected_fragments[role][".pptx"] in slide_text

        loaded_odt = load_odf(io.BytesIO(documents[".odt"].content))
        assert any(
            expected_fragments[role][".odt"] in str(node)
            for node in loaded_odt.getElementsByType(P)
        )


def test_write_bundle_records_labels_and_manifest(tmp_path):
    bundle = generate_all()["calibration"]

    write_bundle(bundle, tmp_path)

    assert (tmp_path / "labels.json").is_file()
    assert (tmp_path / "manifest.json").is_file()
    assert len(list((tmp_path / "documents").rglob("*.*"))) == len(bundle.documents)


def test_write_all_corpora_removes_obsolete_single_corpus_layout(tmp_path):
    legacy_documents = tmp_path / "documents"
    legacy_documents.mkdir()
    (legacy_documents / "old.txt").write_text("obsolete", encoding="utf-8")
    (tmp_path / "labels.json").write_text("[]", encoding="utf-8")
    (tmp_path / "manifest.json").write_text("{}", encoding="utf-8")

    manifests = write_all_corpora(tmp_path)

    assert set(manifests) == {"calibration", "holdout"}
    assert not legacy_documents.exists()
    assert not (tmp_path / "labels.json").exists()
    assert not (tmp_path / "manifest.json").exists()
    assert (tmp_path / "calibration" / "manifest.json").is_file()
    assert (tmp_path / "holdout" / "manifest.json").is_file()
